from pptx import Presentation
from pptx.util import Inches

# Yeni bir sunu oluşturma
prs = Presentation()

# Dosya yolu
img_folder_path = "Images/"

# Dosyadaki tüm resim dosyalarını alma
import os
image_files = [f for f in os.listdir(img_folder_path) if os.path.isfile(os.path.join(img_folder_path, f))]

# Her resim için
for image_file in image_files:
    # Yeni bir slayt oluşturma
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Resmi ekleme
    img_path = os.path.join(img_folder_path, image_file)
    left = Inches(1)
    top = Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width=Inches(4), height=Inches(5))

    # Başlık ekleme (dosya adı kullanıldı şimdilik)
    title_shape = slide.shapes.title
    title_shape.text = os.path.splitext(image_file)[0]

    # Açıklama ekleme (dosya adı kullanıldı şimdilik)
    left = Inches(5)
    top = Inches(2)  # Resmin altında olacak şekilde ayarlandı
    width = Inches(8)
    height = Inches(1)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame

    p = content_frame.add_paragraph()
    p.text = "Açıklama: " + os.path.splitext(image_file)[0]

# Sunuyu dosyaya kaydetme
prs.save("ornek.pptx")